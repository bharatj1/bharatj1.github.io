"""
File content extractor — Word, Excel, PowerPoint, PDF fallback.
Images are handled natively by Claude vision, no extraction needed.
"""
import base64
import io


def extract(filename, data_bytes):
    """
    Extract text content from a document.
    Returns a string of the content.
    """
    name = filename.lower()

    if name.endswith(".docx"):
        return _extract_docx(data_bytes)
    elif name.endswith(".xlsx") or name.endswith(".xls"):
        return _extract_excel(data_bytes)
    elif name.endswith(".pptx") or name.endswith(".ppt"):
        return _extract_pptx(data_bytes)
    elif name.endswith(".pdf"):
        return _extract_pdf(data_bytes)
    elif name.endswith(".txt") or name.endswith(".csv") or name.endswith(".md"):
        return data_bytes.decode("utf-8", errors="replace")
    else:
        return None  # unsupported — will be skipped


def _extract_docx(data):
    try:
        import docx
        doc = docx.Document(io.BytesIO(data))
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text.strip() for c in row.cells))
        return "\n".join(parts)
    except ImportError:
        return "[Install python-docx to read Word files: pip install python-docx]"
    except Exception as e:
        return f"[Could not read Word file: {e}]"


def _extract_excel(data):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
        parts = []
        for sheet in wb.worksheets:
            parts.append(f"=== Sheet: {sheet.title} ===")
            for row in sheet.iter_rows(values_only=True):
                row_vals = [str(v) if v is not None else "" for v in row]
                if any(v.strip() for v in row_vals):
                    parts.append(" | ".join(row_vals))
        return "\n".join(parts)
    except ImportError:
        return "[Install openpyxl to read Excel files: pip install openpyxl]"
    except Exception as e:
        return f"[Could not read Excel file: {e}]"


def _extract_pptx(data):
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"=== Slide {i} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n".join(parts)
    except ImportError:
        return "[Install python-pptx to read PowerPoint files: pip install python-pptx]"
    except Exception as e:
        return f"[Could not read PowerPoint file: {e}]"


def _extract_pdf(data):
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    parts.append(f"=== Page {i} ===\n{text.strip()}")
        return "\n\n".join(parts)
    except ImportError:
        return "[Install pdfplumber to read PDFs: pip install pdfplumber]"
    except Exception as e:
        return f"[Could not read PDF: {e}]"
